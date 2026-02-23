"""
PPT 演示文稿生成服务
将博客内容转化为演示文稿，支持 python-pptx 直接生成 PPTX

参考来源: deer-flow/skills/public/ppt-generation/scripts/generate.py
"""
import logging
import os
from dataclasses import dataclass, field
from enum import Enum
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)


class SlideStyle(Enum):
    """PPT 设计风格（参考 DeerFlow SKILL.md 8 种风格，选取 6 种适配）"""
    GLASSMORPHISM = "glassmorphism"
    DARK_PREMIUM = "dark-premium"
    GRADIENT_MODERN = "gradient-modern"
    MINIMAL_SWISS = "minimal-swiss"
    KEYNOTE = "keynote"
    EDITORIAL = "editorial"


@dataclass
class SlideContent:
    """单张幻灯片内容"""
    slide_number: int
    slide_type: str  # title / content / conclusion
    title: str
    subtitle: str = ""
    key_points: List[str] = field(default_factory=list)
    speaker_notes: str = ""
    image_url: Optional[str] = None


@dataclass
class PresentationPlan:
    """演示文稿计划"""
    title: str
    style: SlideStyle = SlideStyle.KEYNOTE
    aspect_ratio: str = "16:9"
    slides: List[SlideContent] = field(default_factory=list)


class PPTService:
    """PPT 生成服务"""

    STYLE_THEMES: Dict[str, str] = {
        "glassmorphism": "gaia",
        "dark-premium": "uncover",
        "gradient-modern": "gaia",
        "minimal-swiss": "default",
        "keynote": "uncover",
        "editorial": "default",
    }

    def __init__(self, llm_service=None):
        self.llm_service = llm_service
        self.output_folder = os.path.join(
            os.path.dirname(os.path.dirname(__file__)), 'outputs', 'ppts'
        )
        os.makedirs(self.output_folder, exist_ok=True)

    def blog_to_plan(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        style: str = "keynote",
        aspect_ratio: str = "16:9",
    ) -> PresentationPlan:
        """将博客章节结构转化为演示文稿计划"""
        slides: List[SlideContent] = []

        # 封面页
        subtitle = ""
        if sections:
            subtitle = sections[0].get("summary", "")
        slides.append(SlideContent(
            slide_number=1,
            slide_type="title",
            title=title,
            subtitle=subtitle,
        ))

        # 内容页：每个章节一张 Slide
        for i, section in enumerate(sections):
            key_points = section.get("key_points", [])
            if not key_points and section.get("content"):
                paragraphs = section["content"].split("\n\n")
                key_points = [
                    p.split("。")[0] + "。"
                    for p in paragraphs[:3] if p.strip()
                ]
            slides.append(SlideContent(
                slide_number=len(slides) + 1,
                slide_type="content",
                title=section.get("title", f"第 {i+1} 章"),
                key_points=key_points[:5],
                speaker_notes=section.get("content", "")[:500],
            ))

        # 总结页
        slides.append(SlideContent(
            slide_number=len(slides) + 1,
            slide_type="conclusion",
            title="总结",
            subtitle="感谢阅读",
        ))

        return PresentationPlan(
            title=title,
            style=SlideStyle(style),
            aspect_ratio=aspect_ratio,
            slides=slides,
        )

    def plan_to_marp_markdown(self, plan: PresentationPlan) -> str:
        """将演示计划转为 Marp Markdown"""
        theme = self.STYLE_THEMES.get(plan.style.value, "default")
        lines = [
            "---",
            "marp: true",
            f"theme: {theme}",
            f"size: {plan.aspect_ratio}",
            "paginate: true",
            "---",
            "",
        ]

        for slide in plan.slides:
            if slide.slide_type == "title":
                lines.append(f"# {slide.title}")
                lines.append("")
                if slide.subtitle:
                    lines.append(f"### {slide.subtitle}")
                    lines.append("")
                lines.append("---")
                lines.append("")
            elif slide.slide_type == "content":
                lines.append(f"## {slide.title}")
                lines.append("")
                for point in slide.key_points:
                    lines.append(f"- {point}")
                lines.append("")
                lines.append("---")
                lines.append("")
            elif slide.slide_type == "conclusion":
                lines.append(f"# {slide.title}")
                lines.append("")
                if slide.subtitle:
                    lines.append(f"### {slide.subtitle}")
                    lines.append("")

        return "\n".join(lines)

    def generate_pptx(
        self,
        plan: PresentationPlan,
        output_path: str,
    ) -> str:
        """使用 python-pptx 直接生成 PPTX 文件"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()

        # 设置宽高比
        if plan.aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        for slide_content in plan.slides:
            if slide_content.slide_type == "title":
                layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_content.title
                if len(slide.placeholders) > 1 and slide_content.subtitle:
                    slide.placeholders[1].text = slide_content.subtitle
            elif slide_content.slide_type == "content":
                layout = prs.slide_layouts[1]  # Title + Content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_content.title
                if len(slide.placeholders) > 1 and slide_content.key_points:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for j, point in enumerate(slide_content.key_points):
                        if j == 0:
                            tf.paragraphs[0].text = point
                            tf.paragraphs[0].font.size = Pt(18)
                        else:
                            p = tf.add_paragraph()
                            p.text = point
                            p.font.size = Pt(18)
            elif slide_content.slide_type == "conclusion":
                layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_content.title
                if len(slide.placeholders) > 1 and slide_content.subtitle:
                    slide.placeholders[1].text = slide_content.subtitle

            # 添加演讲者备注
            if slide_content.speaker_notes:
                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame
                tf.text = slide_content.speaker_notes

        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        prs.save(output_path)
        return output_path

    def generate(
        self,
        task_id: str,
        title: str,
        sections: List[Dict[str, Any]],
        style: str = "keynote",
        aspect_ratio: str = "16:9",
        task_manager=None,
    ) -> Dict[str, Any]:
        """完整 PPT 生成管线"""
        def _send_progress(stage: str, progress: int, message: str):
            if task_manager:
                task_manager.send_progress(task_id, stage, progress, message)

        try:
            # Stage 1: 生成计划
            _send_progress("plan", 0, "正在生成演示计划...")
            plan = self.blog_to_plan(title, sections, style, aspect_ratio)
            _send_progress("plan", 100,
                           f"计划生成完成: {len(plan.slides)} 张幻灯片")

            # Stage 2: 生成 Marp Markdown
            _send_progress("markdown", 0, "正在生成 Marp Markdown...")
            marp_md = self.plan_to_marp_markdown(plan)
            _send_progress("markdown", 100, "Markdown 生成完成")

            # Stage 3: 生成 PPTX
            _send_progress("compose", 0, "正在合成 PPTX...")
            filename = f"ppt_{task_id}.pptx"
            output_path = os.path.join(self.output_folder, filename)
            self.generate_pptx(plan, output_path)
            _send_progress("compose", 100, "PPT 生成完成")

            return {
                "success": True,
                "pptx_path": output_path,
                "marp_markdown": marp_md,
                "slide_count": len(plan.slides),
                "style": style,
            }
        except Exception as e:
            logger.error(f"PPT 生成失败: {e}", exc_info=True)
            if task_manager:
                task_manager.send_error(task_id, "ppt", str(e))
            return {"success": False, "error": str(e)}


# ── 全局服务实例 ─────────────────────────────────────────────

_ppt_service: Optional[PPTService] = None


def get_ppt_service() -> Optional[PPTService]:
    return _ppt_service


def init_ppt_service(llm_service=None) -> PPTService:
    global _ppt_service
    _ppt_service = PPTService(llm_service=llm_service)
    logger.info("PPT 演示文稿生成服务已初始化")
    return _ppt_service
