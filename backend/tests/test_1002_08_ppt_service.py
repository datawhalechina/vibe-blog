"""
PPT 演示文稿生成服务 — 单元测试
TDD Red Phase: 先写测试，再实现代码

覆盖:
- SlideContent / PresentationPlan / SlideStyle 数据模型
- PPTService.blog_to_plan() 博客→演示计划转换
- PPTService.plan_to_marp_markdown() Marp Markdown 生成
- PPTService.generate_pptx() PPTX 文件生成（python-pptx 回退方案）
- 风格主题映射完整性
- API 路由
"""

import os
import sys
import json
import tempfile
import pytest
from unittest.mock import patch, MagicMock

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))


# ── 数据模型测试 ──────────────────────────────────────────────

class TestSlideStyle:
    """SlideStyle 枚举"""

    def test_all_styles_defined(self):
        from services.ppt_service import SlideStyle
        expected = {"glassmorphism", "dark-premium", "gradient-modern",
                    "minimal-swiss", "keynote", "editorial"}
        actual = {s.value for s in SlideStyle}
        assert expected == actual

    def test_style_from_string(self):
        from services.ppt_service import SlideStyle
        assert SlideStyle("keynote") == SlideStyle.KEYNOTE
        assert SlideStyle("dark-premium") == SlideStyle.DARK_PREMIUM


class TestSlideContent:
    """SlideContent 数据模型"""

    def test_create_slide_content(self):
        from services.ppt_service import SlideContent
        slide = SlideContent(
            slide_number=1,
            slide_type="title",
            title="Hello World",
        )
        assert slide.slide_number == 1
        assert slide.slide_type == "title"
        assert slide.title == "Hello World"
        assert slide.key_points == []
        assert slide.speaker_notes == ""
    def test_slide_content_with_key_points(self):
        from services.ppt_service import SlideContent
        slide = SlideContent(
            slide_number=2,
            slide_type="content",
            title="核心概念",
            key_points=["StateGraph", "Node", "Edge"],
            speaker_notes="详细说明...",
        )
        assert len(slide.key_points) == 3
        assert slide.speaker_notes == "详细说明..."


class TestPresentationPlan:
    """PresentationPlan 数据模型"""

    def test_create_plan_defaults(self):
        from services.ppt_service import PresentationPlan, SlideStyle
        plan = PresentationPlan(title="Test")
        assert plan.title == "Test"
        assert plan.style == SlideStyle.KEYNOTE
        assert plan.aspect_ratio == "16:9"
        assert plan.slides == []

    def test_create_plan_with_slides(self):
        from services.ppt_service import PresentationPlan, SlideContent, SlideStyle
        slides = [
            SlideContent(slide_number=1, slide_type="title", title="封面"),
            SlideContent(slide_number=2, slide_type="content", title="内容"),
        ]
        plan = PresentationPlan(
            title="My PPT",
            style=SlideStyle.DARK_PREMIUM,
            aspect_ratio="4:3",
            slides=slides,
        )
        assert len(plan.slides) == 2
        assert plan.style == SlideStyle.DARK_PREMIUM
        assert plan.aspect_ratio == "4:3"


# ── PPTService 核心逻辑测试 ──────────────────────────────────

class TestBlogToPlan:
    """博客内容 → 演示计划转换"""

    def setup_method(self):
        from services.ppt_service import PPTService
        self.service = PPTService()
        self.sample_sections = [
            {
                "title": "什么是 LangGraph",
                "content": "LangGraph 是一个基于图的编排框架。它支持状态管理。",
                "key_points": ["图编排", "状态管理", "条件路由"],
            },
            {
                "title": "核心概念",
                "content": "StateGraph 是核心类。Node 是执行单元。Edge 定义流转。",
                "key_points": ["StateGraph", "Node", "Edge"],
            },
            {
                "title": "实战案例",
                "content": "构建一个多 Agent 系统。",
                "key_points": ["多 Agent", "工具调用"],
            },
        ]

    def test_blog_to_plan_basic(self):
        """3 章博客应生成 5 张 Slide（封面+3内容+总结）"""
        plan = self.service.blog_to_plan(
            title="LangGraph 入门指南",
            sections=self.sample_sections,
        )
        assert len(plan.slides) == 5
        assert plan.slides[0].slide_type == "title"
        assert plan.slides[-1].slide_type == "conclusion"
        assert all(s.slide_type == "content" for s in plan.slides[1:4])

    def test_blog_to_plan_title_slide(self):
        """封面页标题应与博客标题一致"""
        plan = self.service.blog_to_plan(
            title="LangGraph 入门指南",
            sections=self.sample_sections,
        )
        assert plan.slides[0].title == "LangGraph 入门指南"
        assert plan.title == "LangGraph 入门指南"

    def test_blog_to_plan_empty_sections(self):
        """空章节列表应只生成封面+总结"""
        plan = self.service.blog_to_plan(title="空文章", sections=[])
        assert len(plan.slides) == 2
        assert plan.slides[0].slide_type == "title"
        assert plan.slides[-1].slide_type == "conclusion"

    def test_key_points_limit(self):
        """每页最多 5 个要点"""
        sections = [{
            "title": "测试",
            "key_points": [f"要点{i}" for i in range(10)],
        }]
        plan = self.service.blog_to_plan(title="测试", sections=sections)
        content_slides = [s for s in plan.slides if s.slide_type == "content"]
        assert len(content_slides[0].key_points) <= 5

    def test_blog_to_plan_extracts_points_from_content(self):
        """无 key_points 时从 content 提取"""
        sections = [{
            "title": "测试章节",
            "content": "第一段内容。详细说明。\n\n第二段内容。更多细节。\n\n第三段内容。总结。",
        }]
        plan = self.service.blog_to_plan(title="测试", sections=sections)
        content_slide = [s for s in plan.slides if s.slide_type == "content"][0]
        assert len(content_slide.key_points) > 0

    def test_blog_to_plan_style_param(self):
        """风格参数应正确传递"""
        plan = self.service.blog_to_plan(
            title="测试",
            sections=self.sample_sections,
            style="dark-premium",
        )
        from services.ppt_service import SlideStyle
        assert plan.style == SlideStyle.DARK_PREMIUM

    def test_blog_to_plan_aspect_ratio(self):
        """宽高比参数应正确传递"""
        plan = self.service.blog_to_plan(
            title="测试",
            sections=self.sample_sections,
            aspect_ratio="4:3",
        )
        assert plan.aspect_ratio == "4:3"


class TestPlanToMarpMarkdown:
    """演示计划 → Marp Markdown 转换"""

    def setup_method(self):
        from services.ppt_service import PPTService, PresentationPlan, SlideContent, SlideStyle
        self.service = PPTService()
        self.plan = PresentationPlan(
            title="测试演示",
            style=SlideStyle.KEYNOTE,
            aspect_ratio="16:9",
            slides=[
                SlideContent(slide_number=1, slide_type="title",
                             title="测试演示", subtitle="副标题"),
                SlideContent(slide_number=2, slide_type="content",
                             title="第一章", key_points=["要点A", "要点B"]),
                SlideContent(slide_number=3, slide_type="conclusion",
                             title="总结", subtitle="感谢阅读"),
            ],
        )

    def test_marp_frontmatter(self):
        """Marp Markdown 应包含正确的 frontmatter"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        assert marp.startswith("---\nmarp: true")
        assert "theme: uncover" in marp  # keynote → uncover
        assert "paginate: true" in marp

    def test_marp_contains_separators(self):
        """Marp Markdown 应包含 --- 分隔符"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        # frontmatter 的 --- 不算，内容分隔符至少有 2 个（title→content, content→conclusion）
        content_after_frontmatter = marp.split("---\n", 2)[-1]
        assert "\n---\n" in content_after_frontmatter

    def test_marp_title_slide(self):
        """封面页应包含 # 标题"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        assert "# 测试演示" in marp
        assert "### 副标题" in marp

    def test_marp_content_slide(self):
        """内容页应包含 ## 标题和要点列表"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        assert "## 第一章" in marp
        assert "- 要点A" in marp
        assert "- 要点B" in marp

    def test_marp_conclusion_slide(self):
        """总结页应包含 # 标题"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        assert "# 总结" in marp

    def test_plan_to_marp_keynote_style(self):
        """keynote 风格 → Marp theme: uncover"""
        marp = self.service.plan_to_marp_markdown(self.plan)
        assert "theme: uncover" in marp

    def test_plan_to_marp_glassmorphism_style(self):
        """glassmorphism 风格 → Marp theme: gaia"""
        from services.ppt_service import PresentationPlan, SlideContent, SlideStyle
        plan = PresentationPlan(
            title="Test",
            style=SlideStyle.GLASSMORPHISM,
            slides=[SlideContent(slide_number=1, slide_type="title", title="Test")],
        )
        marp = self.service.plan_to_marp_markdown(plan)
        assert "theme: gaia" in marp


class TestStyleThemeMapping:
    """风格主题映射完整性"""

    def test_all_styles_mapped(self):
        """所有 SlideStyle 枚举值都应映射到有效 Marp 主题"""
        from services.ppt_service import PPTService, SlideStyle
        for style in SlideStyle:
            assert style.value in PPTService.STYLE_THEMES, \
                f"风格 {style.value} 未映射到 Marp 主题"

    def test_valid_marp_themes(self):
        """映射的主题应是有效的 Marp 内置主题"""
        from services.ppt_service import PPTService
        valid_themes = {"default", "gaia", "uncover"}
        for style, theme in PPTService.STYLE_THEMES.items():
            assert theme in valid_themes, \
                f"风格 {style} 映射到无效主题 {theme}"


# ── PPTX 生成测试（python-pptx 回退方案）────────────────────

class TestGeneratePptxFallback:
    """python-pptx 直接生成 PPTX（无需 Marp CLI）"""

    def setup_method(self):
        from services.ppt_service import PPTService
        self.service = PPTService()

    def test_generate_pptx_creates_file(self):
        """generate_pptx 应生成有效的 PPTX 文件"""
        from services.ppt_service import PresentationPlan, SlideContent, SlideStyle
        plan = PresentationPlan(
            title="测试 PPT",
            style=SlideStyle.KEYNOTE,
            slides=[
                SlideContent(slide_number=1, slide_type="title",
                             title="测试标题", subtitle="测试副标题"),
                SlideContent(slide_number=2, slide_type="content",
                             title="内容页",
                             key_points=["要点1", "要点2", "要点3"]),
                SlideContent(slide_number=3, slide_type="conclusion",
                             title="总结", subtitle="谢谢"),
            ],
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name
        try:
            result = self.service.generate_pptx(plan, output_path)
            assert result == output_path
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0
        finally:
            os.unlink(output_path)

    def test_generate_pptx_16_9(self):
        """16:9 宽高比 PPTX 文件可正常生成"""
        from services.ppt_service import PresentationPlan, SlideContent
        plan = PresentationPlan(
            title="16:9 Test",
            aspect_ratio="16:9",
            slides=[SlideContent(slide_number=1, slide_type="title", title="Test")],
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name
        try:
            self.service.generate_pptx(plan, output_path)
            from pptx import Presentation
            prs = Presentation(output_path)
            from pptx.util import Inches
            assert prs.slide_width == Inches(13.333)
            assert prs.slide_height == Inches(7.5)
        finally:
            os.unlink(output_path)

    def test_generate_pptx_4_3(self):
        """4:3 宽高比 PPTX 文件可正常生成"""
        from services.ppt_service import PresentationPlan, SlideContent
        plan = PresentationPlan(
            title="4:3 Test",
            aspect_ratio="4:3",
            slides=[SlideContent(slide_number=1, slide_type="title", title="Test")],
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name
        try:
            self.service.generate_pptx(plan, output_path)
            from pptx import Presentation
            prs = Presentation(output_path)
            from pptx.util import Inches
            assert prs.slide_width == Inches(10)
            assert prs.slide_height == Inches(7.5)
        finally:
            os.unlink(output_path)

    def test_generate_pptx_slide_count(self):
        """生成的 PPTX 应包含正确数量的 Slide"""
        from services.ppt_service import PresentationPlan, SlideContent
        plan = PresentationPlan(
            title="Count Test",
            slides=[
                SlideContent(slide_number=i, slide_type="content",
                             title=f"Slide {i}", key_points=[f"Point {i}"])
                for i in range(1, 6)
            ],
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name
        try:
            self.service.generate_pptx(plan, output_path)
            from pptx import Presentation
            prs = Presentation(output_path)
            assert len(prs.slides) == 5
        finally:
            os.unlink(output_path)

    def test_generate_pptx_speaker_notes(self):
        """PPTX 应包含演讲者备注"""
        from services.ppt_service import PresentationPlan, SlideContent
        plan = PresentationPlan(
            title="Notes Test",
            slides=[
                SlideContent(slide_number=1, slide_type="content",
                             title="有备注的页",
                             key_points=["要点1"],
                             speaker_notes="这是演讲者备注内容"),
            ],
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name
        try:
            self.service.generate_pptx(plan, output_path)
            from pptx import Presentation
            prs = Presentation(output_path)
            slide = prs.slides[0]
            notes_text = slide.notes_slide.notes_text_frame.text
            assert "这是演讲者备注内容" in notes_text
        finally:
            os.unlink(output_path)


# ── 全局服务实例测试 ─────────────────────────────────────────

class TestPPTServiceSingleton:
    """全局服务实例管理"""

    def test_init_and_get(self):
        from services.ppt_service import init_ppt_service, get_ppt_service
        service = init_ppt_service()
        assert service is not None
        assert get_ppt_service() is service

    def test_get_before_init_returns_none(self):
        import services.ppt_service as mod
        mod._ppt_service = None
        assert mod.get_ppt_service() is None


# ── API 路由测试 ──────────────────────────────────────────────

class TestPPTRoutes:
    """PPT API 路由"""

    @pytest.fixture
    def app(self):
        from flask import Flask
        from routes.ppt_routes import ppt_bp
        app = Flask(__name__)
        app.register_blueprint(ppt_bp)
        app.config['TESTING'] = True
        return app

    @pytest.fixture
    def client(self, app):
        return app.test_client()

    def test_get_styles(self, client):
        """GET /api/ppt/styles 返回风格列表"""
        resp = client.get('/api/ppt/styles')
        assert resp.status_code == 200
        data = resp.get_json()
        assert data['success'] is True
        assert len(data['styles']) == 6

    def test_generate_missing_content(self, client):
        """缺少 content 参数返回 400"""
        resp = client.post('/api/ppt/generate', json={})
        assert resp.status_code == 400
        data = resp.get_json()
        assert data['success'] is False

    @patch('routes.ppt_routes.get_ppt_service')
    def test_generate_service_unavailable(self, mock_get_svc, client):
        """服务不可用返回 503"""
        mock_get_svc.return_value = None
        resp = client.post('/api/ppt/generate', json={
            'content': 'test',
            'title': 'test',
            'sections': [{"title": "ch1", "key_points": ["p1"]}],
        })
        assert resp.status_code == 503

    @patch('routes.ppt_routes.get_ppt_service')
    @patch('routes.ppt_routes.get_task_manager')
    def test_generate_success(self, mock_get_tm, mock_get_svc, client):
        """正常生成返回 task_id"""
        mock_svc = MagicMock()
        mock_get_svc.return_value = mock_svc

        mock_tm = MagicMock()
        mock_tm.create_task.return_value = 'test-task-id'
        mock_get_tm.return_value = mock_tm

        resp = client.post('/api/ppt/generate', json={
            'content': '一篇关于 AI 的文章',
            'title': 'AI PPT',
            'sections': [{"title": "ch1", "key_points": ["p1"]}],
        })
        assert resp.status_code == 200
        data = resp.get_json()
        assert data['success'] is True
        assert 'task_id' in data

    def test_download_not_found(self, client):
        """下载不存在的任务返回 404"""
        resp = client.get('/api/ppt/download/nonexistent')
        assert resp.status_code == 404

